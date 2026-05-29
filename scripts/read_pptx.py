#!/usr/bin/env python3
"""Extract text + speaker notes from a PowerPoint .pptx.

Usage:
    read_pptx.py deck.pptx               # slides as markdown
    read_pptx.py deck.pptx --json        # structured JSON
    read_pptx.py deck.pptx --notes-only  # speaker notes only
"""
import argparse
import json
import sys


def main():
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("pptx")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--notes-only", action="store_true")
    a = ap.parse_args()

    from pptx import Presentation
    prs = Presentation(a.pptx)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs).strip()
                    if line:
                        texts.append(line)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"slide": i, "text": texts, "notes": notes})

    if a.json:
        print(json.dumps(slides, indent=2))
    elif a.notes_only:
        for s in slides:
            if s["notes"]:
                print(f"## Slide {s['slide']} notes\n{s['notes']}\n")
    else:
        for s in slides:
            print(f"## Slide {s['slide']}")
            for t in s["text"]:
                print(f"- {t}")
            if s["notes"]:
                print(f"\n*Notes:* {s['notes']}")
            print()


if __name__ == "__main__":
    try:
        main()
    except FileNotFoundError as e:
        print(f"ERROR: {e}", file=sys.stderr); sys.exit(2)
